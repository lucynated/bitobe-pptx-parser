"""PPTX parser microservice for n8n workflow."""
from fastapi import FastAPI, UploadFile, File, Request
from fastapi.responses import JSONResponse
from pptx import Presentation
import httpx
import io

app = FastAPI()


def parse_presentation(content: bytes) -> dict:
    """Parse PPTX bytes into structured slide data."""
    prs = Presentation(io.BytesIO(content))
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        texts.append(t)
            if shape.has_table:
                for row in shape.table.rows:
                    row_texts = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if row_texts:
                        texts.append("ТАБЛИЦА: " + " | ".join(row_texts))
        notes = ""
        if slide.has_notes_slide:
            for para in slide.notes_slide.notes_text_frame.paragraphs:
                nt = para.text.strip()
                if nt and not nt.isdigit():
                    notes += nt + "\n"
        slides.append({
            "number": i,
            "title": texts[0] if texts else f"Слайд {i}",
            "content": "\n".join(texts) or "(пустой слайд)",
            "notes": notes.strip(),
            "layout_type": "title" if i == 1 else "content",
        })
    return {"slides": slides, "total_slides": len(slides), "file_type": "pptx"}


@app.post("/parse")
async def parse_pptx(file: UploadFile = File(...)):
    """Parse uploaded PPTX file."""
    try:
        content = await file.read()
        return parse_presentation(content)
    except Exception as e:
        return JSONResponse(status_code=400, content={"error": str(e)})


@app.post("/parse-url")
async def parse_pptx_url(request: Request):
    """Download PPTX from URL and parse."""
    try:
        data = await request.json()
        file_url = data.get("file_url")
        if not file_url:
            return JSONResponse(status_code=400, content={"error": "file_url обязателен"})
        async with httpx.AsyncClient(follow_redirects=True, timeout=30) as client:
            resp = await client.get(file_url)
            resp.raise_for_status()
        return parse_presentation(resp.content)
    except httpx.HTTPError as e:
        return JSONResponse(status_code=502, content={"error": f"Ошибка загрузки: {e}"})
    except Exception as e:
        return JSONResponse(status_code=400, content={"error": str(e)})


@app.get("/health")
def health():
    return {"status": "ok"}
